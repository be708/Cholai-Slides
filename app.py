import os
from flask import Flask, request, render_template, send_from_directory
from flask_cors import CORS 
from pptx import Presentation
from pptx.util import Inches, Pt
import google.generativeai as genai
from io import BytesIO

app = Flask(__name__)
CORS(app)
GOOGLE_API_KEY = os.environ.get("GOOGLE_API_KEY")
genai.configure(api_key=GOOGLE_API_KEY)

TEMPLATES = {
    "1": "1.png",
    "2": "2.png", 
    "3": "3.png",
    #... add up to 27
    "27": "27.png"
}


def generate_with_ai(user_idea):
    model = genai.GenerativeModel('gemini-1.0-pro-latest')
    prompt = f"""
    You are an expert presentation designer. Create a 6 slide PowerPoint outline for: {user_idea}
    Format: Use Markdown. Each slide starts with # Title
    Bullet points start with - 
    Separate slides with ---
    Example:
    # Problem
    - Point 1
    - Point 2
    ---
    # Solution
    - Point 1
    """
    response = model.generate_content(prompt)
    return parse_markdown(response.text)


def parse_markdown(md_text):
    slides = []
    blocks = md_text.split('---')
    for block in blocks:
        lines = [l.strip() for l in block.split('\n') if l.strip()]
        if not lines: 
            continue
        title = lines[0].replace('# ', '')
        content = [l.replace('- ', '') for l in lines[1:]]
        slides.append({"title": title, "content": content})
    return slides


@app.route('/')
def home():
    return render_template('index.html')

@app.route('/generate', methods=['POST'])
def generate_slides():
    try:
        user_idea = request.form.get('topic')
        manual_markdown = request.form.get('markdown')
        email = request.form.get('email', '')
        template_choice = request.form.get('template', '1')

        if user_idea and user_idea.strip():
            slides_data = generate_with_ai(user_idea)
        else:
            slides_data = parse_markdown(manual_markdown)

        product_name = user_idea[:30] if user_idea else "Manual_Slides"
        template_file = TEMPLATES.get(template_choice)
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        for i, slide_info in enumerate(slides_data):
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # 1. ADD PNG AS BACKGROUND
            template_num = str((i % 27) + 1)
            template_png = TEMPLATES.get(template_num)
            slide.shapes.add_picture(f'templates/{template_png}', 0, 0, width=prs.slide_width, height=prs.slide_height)

            # 2. ADD TEXT BOX ON TOP OF PNG
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.clear()

            pe = tf.add_paragraph()
            pe.text = slide_info["title"]
            pe.font.size = Pt(32)
            pe.font.bold = True

            for point in slide_info["content"]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                p.font.size = Pt(18)

        file_stream = BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)

        return send_file(
            file_stream,
            as_attachment=True,
            download_name=f"{product_name}.pptx",
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception as e:
        return f"<h1>Error: {str(e)}</h1>", 500
        if__name__=='__main__':
           app.run(debug=True)
